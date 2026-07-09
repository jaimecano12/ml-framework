"""Build presentation.pptx — 25-minute defense slides (English, PowerPoint).

Native, fully editable PPTX with a professional design system based on the
UPM visual identity (upmBlue #243F60, upmOrange #FF8000):

  * full-bleed title and closing slides
  * agenda + five section dividers with progress dots
  * consistent header bar, footer and slide numbering on content slides
  * KPI cards, numbered contribution cards, styled tables, native charts

Usage:  python scripts/build_presentation_pptx.py
Output: presentation.pptx (repository root)
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---- palette ----------------------------------------------------------
UPM_BLUE = RGBColor(0x24, 0x3F, 0x60)
UPM_BLUE_DK = RGBColor(0x1A, 0x2F, 0x49)
UPM_ORANGE = RGBColor(0xFF, 0x80, 0x00)
TEAL = RGBColor(0x2E, 0x8B, 0x8B)
TEAL_LIGHT = RGBColor(0xDD, 0xEE, 0xEE)
BLUE_LIGHT = RGBColor(0xE8, 0xEF, 0xF9)
BLUE_PALE = RGBColor(0xF3, 0xF6, 0xFB)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE6)
ALT_ROW = RGBColor(0xED, 0xF1, 0xF7)
GRAY = RGBColor(0x8A, 0x8A, 0x8A)
GRAY_LIGHT = RGBColor(0xD5, 0xD9, 0xE0)
GRAY_TXT = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_SOFT = RGBColor(0xB9, 0xC8, 0xDC)      # light text on blue bg
GREEN = RGBColor(0x1E, 0x7A, 0x1E)
GREEN_LIGHT = RGBColor(0xD6, 0xEC, 0xD6)
RED = RGBColor(0xB4, 0x1E, 0x1E)
RED_LIGHT = RGBColor(0xFB, 0xEA, 0xEA)
CODE_BG = RGBColor(0xF7, 0xF7, 0xF7)
CODE_FG = RGBColor(0x1E, 0x3A, 0x66)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SECTIONS = [
    "Motivation & Objectives",
    "System & Methodology",
    "Experimental Results",
    "Benchmark",
    "Demo & Conclusions",
]

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

_slide_no = 0  # running number printed in the footer chip


# ---- low-level helpers -------------------------------------------------
def new_slide():
    return prs.slides.add_slide(BLANK)


def set_fill(shape, color, line_color=None, line_w=0.75):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False


def rich(par, text, size=16, color=DARK, base_bold=False, align=PP_ALIGN.LEFT,
         mono=False, font=None):
    """Add runs to a paragraph; **segments** become bold."""
    par.alignment = align
    for i, part in enumerate(text.split("**")):
        if not part:
            continue
        run = par.add_run()
        run.text = part
        run.font.size = Pt(size)
        run.font.bold = base_bold or (i % 2 == 1)
        run.font.color.rgb = color
        if mono:
            run.font.name = "Consolas"
        elif font:
            run.font.name = font
    return par


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    return box


def bullets(slide, items, x, y, w, h, size=16, gap=8, bullet=True, color=DARK):
    box = textbox(slide, x, y, w, h)
    tf = box.text_frame
    for i, item in enumerate(items):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rich(par, ("•  " if bullet else "") + item, size=size, color=color)
        par.space_after = Pt(gap)
    return box


def rounded(slide, x, y, w, h, fill, line=None, line_w=0.75, radius=0.12):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    set_fill(shp, fill, line_color=line, line_w=line_w)
    return shp


def num_circle(slide, n, x, y, d=Inches(0.5), fill=UPM_ORANGE, fg=WHITE, size=18):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    set_fill(c, fill)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(tf.paragraphs[0], str(n), size=size, color=fg, base_bold=True,
         align=PP_ALIGN.CENTER)
    return c


def arrow(slide, x1, y1, x2, y2, color=GRAY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.2)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


# ---- design-system components ------------------------------------------
def chrome(slide, title, kicker=None):
    """Header bar + orange accent + footer with slide number."""
    global _slide_no
    _slide_no += 1
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    set_fill(bar, UPM_BLUE)
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kicker:
        rich(tf.paragraphs[0], kicker.upper(), size=11, color=UPM_ORANGE,
             base_bold=True)
        p = tf.add_paragraph()
        rich(p, title, size=24, color=WHITE, base_bold=True)
    else:
        rich(tf.paragraphs[0], title, size=25, color=WHITE, base_bold=True)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SLIDE_W, Inches(0.04))
    set_fill(accent, UPM_ORANGE)
    # footer
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.12), SLIDE_W - Inches(1.0),
        Pt(1))
    set_fill(line, GRAY_LIGHT)
    ft = textbox(slide, Inches(0.5), Inches(7.16), Inches(9.0), Inches(0.3))
    rich(ft.text_frame.paragraphs[0],
         "Dataset Quality Assessment & Data Leakage Detection  ·  J. Cano Moraño",
         size=10, color=GRAY)
    chip = rounded(slide, SLIDE_W - Inches(0.95), Inches(7.14), Inches(0.45),
                   Inches(0.28), UPM_BLUE, radius=0.5)
    chip.text_frame.margin_left = chip.text_frame.margin_top = 0
    chip.text_frame.margin_right = chip.text_frame.margin_bottom = 0
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(chip.text_frame.paragraphs[0], str(_slide_no), size=11, color=WHITE,
         base_bold=True, align=PP_ALIGN.CENTER)
    return slide


def divider(section_idx, subtitle_lines):
    """Full-bleed section divider with number, progress dots, contents."""
    global _slide_no
    _slide_no += 1
    s = new_slide()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, UPM_BLUE)
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    set_fill(side, UPM_ORANGE)
    num = textbox(s, Inches(0.9), Inches(1.45), Inches(4.0), Inches(2.0))
    rich(num.text_frame.paragraphs[0], f"0{section_idx + 1}", size=110,
         color=UPM_ORANGE, base_bold=True)
    t = textbox(s, Inches(0.95), Inches(3.55), Inches(11.5), Inches(1.0))
    rich(t.text_frame.paragraphs[0], SECTIONS[section_idx], size=40,
         color=WHITE, base_bold=True)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.55),
                              Inches(3.2), Inches(0.035))
    set_fill(rule, UPM_ORANGE)
    sub = textbox(s, Inches(1.0), Inches(4.85), Inches(10.8), Inches(1.6))
    tf = sub.text_frame
    for i, line in enumerate(subtitle_lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rich(par, line, size=16, color=BLUE_SOFT)
        par.space_after = Pt(4)
    # progress dots
    dx = SLIDE_W - Inches(2.9)
    for i, name in enumerate(SECTIONS):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, dx + Inches(i * 0.42),
                                 Inches(0.65), Inches(0.22), Inches(0.22))
        set_fill(dot, UPM_ORANGE if i == section_idx else RGBColor(0x4A, 0x63, 0x84))
    return s


def block(slide, header, body, x, y, w, h, head_color=UPM_ORANGE,
          body_fill=ORANGE_LIGHT, size=15, head_size=14,
          body_align=PP_ALIGN.LEFT):
    hd = rounded(slide, x, y, w, Inches(0.4), head_color, radius=0.3)
    hd.text_frame.margin_left = Inches(0.15)
    hd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(hd.text_frame.paragraphs[0], header, size=head_size, color=WHITE,
         base_bold=True)
    bd = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y + Inches(0.4), w, h - Inches(0.4))
    set_fill(bd, body_fill, line_color=head_color, line_w=0.5)
    tf = bd.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rich(par, line, size=size, align=body_align)
        par.space_after = Pt(5)
    return bd


def kpi_card(slide, value, label, x, y, w=Inches(2.6), h=Inches(1.35),
             accent=UPM_ORANGE):
    card = rounded(slide, x, y, w, h, WHITE, line=GRAY_LIGHT, line_w=1)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.07))
    set_fill(top, accent)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(tf.paragraphs[0], value, size=30, color=UPM_BLUE, base_bold=True,
         align=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    rich(p, label, size=12, color=GRAY_TXT, align=PP_ALIGN.CENTER)
    return card


def numbered_card(slide, n, title, body, x, y, w, h, accent=UPM_ORANGE):
    card = rounded(slide, x, y, w, h, BLUE_PALE, line=GRAY_LIGHT, line_w=1)
    num_circle(slide, n, x + Inches(0.18), y + Inches(0.18), d=Inches(0.44),
               fill=accent, size=16)
    box = textbox(slide, x + Inches(0.75), y + Inches(0.14), w - Inches(0.95),
                  h - Inches(0.3))
    tf = box.text_frame
    rich(tf.paragraphs[0], title, size=15, color=UPM_BLUE, base_bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    rich(p, body, size=12.5, color=DARK)
    return card


def box_node(slide, text, x, y, w, h, fill, line, size=12):
    shp = rounded(slide, x, y, w, h, fill, line=line, line_w=1, radius=0.18)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    rich(tf.paragraphs[0], text, size=size, align=PP_ALIGN.CENTER)
    return shp


def styled_table(slide, data, x, y, w, row_h=0.42, size=13, col_widths=None,
                 align_center_cols=()):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, Inches(row_h * rows))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                UPM_BLUE if r == 0 else (ALT_ROW if r % 2 == 0 else WHITE))
            par = cell.text_frame.paragraphs[0]
            align = PP_ALIGN.CENTER if c in align_center_cols else PP_ALIGN.LEFT
            if val in ("✓", "✗"):
                par.alignment = align
                run = par.add_run()
                run.text = val
                run.font.size = Pt(size)
                run.font.bold = True
                run.font.color.rgb = GREEN if val == "✓" else RED
            else:
                rich(par, val, size=size, color=WHITE if r == 0 else DARK,
                     base_bold=(r == 0), align=align)
    return tbl


def note(slide, text, x=Inches(0.5), y=Inches(6.55), w=None, size=12):
    if w is None:
        w = SLIDE_W - Inches(1.0)
    box = textbox(slide, x, y, w, Inches(0.55))
    rich(box.text_frame.paragraphs[0], text, size=size, color=GRAY_TXT)
    return box


def code_box(slide, lines, x, y, w, h, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, CODE_BG, line_color=RGBColor(0xBB, 0xBB, 0xBB), line_w=0.75)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    set_fill(strip, UPM_ORANGE)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = PP_ALIGN.LEFT
        run = par.add_run()
        run.text = line if line else " "
        run.font.size = Pt(size)
        run.font.name = "Consolas"
        run.font.color.rgb = CODE_FG
    return shp


# ======================================================================
# 1 — Title (full-bleed)
# ======================================================================
_slide_no += 1
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
set_fill(bg, UPM_BLUE)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), SLIDE_W, Inches(1.6))
set_fill(band, UPM_BLUE_DK)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), SLIDE_W, Inches(0.045))
set_fill(accent, UPM_ORANGE)
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
set_fill(side, UPM_ORANGE)

kick = textbox(s, Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.5))
rich(kick.text_frame.paragraphs[0], "MASTER'S THESIS DEFENSE  ·  JULY 2026",
     size=14, color=UPM_ORANGE, base_bold=True)
t = textbox(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(2.4))
tf = t.text_frame
rich(tf.paragraphs[0],
     "An Automated Framework for Dataset Quality Assessment "
     "and Data Leakage Detection in Machine Learning",
     size=34, color=WHITE, base_bold=True)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(4.05),
                          Inches(3.2), Inches(0.035))
set_fill(rule, UPM_ORANGE)
sub = textbox(s, Inches(1.0), Inches(4.35), Inches(11.3), Inches(1.2))
tf = sub.text_frame
rich(tf.paragraphs[0],
     "20 automated checks · unified leakage risk score · LLM semantic analysis "
     "· quantitative benchmark", size=16, color=BLUE_SOFT)
auth = textbox(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(1.1))
tf = auth.text_frame
rich(tf.paragraphs[0], "Jaime Cano Moraño", size=20, color=WHITE, base_bold=True)
p = tf.add_paragraph()
rich(p, "ETSIT — Universidad Politécnica de Madrid   ·   "
        "Illinois Institute of Technology", size=13, color=BLUE_SOFT)

# ======================================================================
# 2 — Agenda
# ======================================================================
s = new_slide()
chrome(s, "Agenda")
agenda = [
    ("Motivation & Objectives", "Why data leakage matters, the gap in existing "
     "tools, and the four contributions", "≈ 5 min"),
    ("System & Methodology", "Architecture, the 20-check catalog, the unified "
     "Leakage Risk Score, LLM semantic analysis, readiness scoring", "≈ 8 min"),
    ("Experimental Results", "12 datasets, readiness scores, the Titanic case "
     "study, impact analysis", "≈ 5 min"),
    ("Benchmark", "Quantitative comparison against ydata-profiling, Deepchecks, "
     "and Great Expectations", "≈ 3 min"),
    ("Demo & Conclusions", "Live Streamlit demo, conclusions, future work",
     "≈ 4 min"),
]
ay = Inches(1.35)
for i, (name, desc, mins) in enumerate(agenda):
    card = rounded(s, Inches(0.7), ay, Inches(11.9), Inches(1.0), BLUE_PALE,
                   line=GRAY_LIGHT, line_w=0.75)
    num_circle(s, i + 1, Inches(0.95), ay + Inches(0.26), d=Inches(0.48))
    box = textbox(s, Inches(1.7), ay + Inches(0.12), Inches(9.3), Inches(0.85))
    tf = box.text_frame
    rich(tf.paragraphs[0], name, size=16, color=UPM_BLUE, base_bold=True)
    p = tf.add_paragraph()
    rich(p, desc, size=12, color=GRAY_TXT)
    mbox = textbox(s, Inches(11.1), ay + Inches(0.3), Inches(1.3), Inches(0.4))
    rich(mbox.text_frame.paragraphs[0], mins, size=13, color=UPM_ORANGE,
         base_bold=True, align=PP_ALIGN.RIGHT)
    ay += Inches(1.13)

# ======================================================================
# DIVIDER 01 — Motivation & Objectives
# ======================================================================
divider(0, [
    "The silent failure mode of machine learning  ·  a motivating example",
    "What existing tools miss  ·  objectives and contributions of this thesis",
])

# ======================================================================
# 3 — The problem
# ======================================================================
s = new_slide()
chrome(s, "Models Fail Because of Data, Not Algorithms",
       kicker="Motivation & Objectives")
bullets(s, [
    "Most ML effort targets **models**; most production failures trace back to "
    "**data**: missing values, duplicates, imbalance, drift … and, most "
    "insidiously, **data leakage**.",
    "Leakage is dangerous because it is silent: the model looks **better**, not "
    "worse — inflated validation metrics that collapse in production.",
    "Documented across medicine, finance, and academic reproducibility studies "
    "(Kaufman et al. 2012; Kapoor & Narayanan).",
], Inches(0.7), Inches(1.35), Inches(7.3), Inches(3.6), size=16, gap=14)
defn = rounded(s, Inches(8.4), Inches(1.4), Inches(4.3), Inches(2.6), BLUE_LIGHT,
               line=UPM_BLUE, line_w=1)
tf = defn.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_right = Inches(0.2)
tf.margin_top = Inches(0.15)
rich(tf.paragraphs[0], "DEFINITION", size=11, color=UPM_ORANGE, base_bold=True)
p = tf.add_paragraph()
p.space_before = Pt(6)
rich(p, "**Data leakage**: information about the target reaches the training "
        "features in a way that will **not exist at prediction time**.",
     size=15, color=UPM_BLUE)
block(s, "Core question of this thesis",
      "Can dataset quality assessment and leakage detection be **automated, "
      "quantified, and made actionable** in a single tool?",
      Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.4), size=17)

# ======================================================================
# 4 — Motivating example
# ======================================================================
s = new_slide()
chrome(s, "A Motivating Example", kicker="Motivation & Objectives")
intro = textbox(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5))
rich(intro.text_frame.paragraphs[0],
     "A logistic regression classifier trained on a dataset with an engineered "
     "leaky feature:", size=17)
block(s, "With leakage (as delivered)", ["**Validation accuracy:  1.000**"],
      Inches(0.9), Inches(2.0), Inches(5.2), Inches(1.35),
      head_color=RED, body_fill=RED_LIGHT, size=22, body_align=PP_ALIGN.CENTER)
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.25), Inches(2.5),
                        Inches(0.85), Inches(0.45))
set_fill(ar, UPM_ORANGE)
block(s, "After removing leaky features", ["**Validation accuracy:  0.952**"],
      Inches(7.25), Inches(2.0), Inches(5.2), Inches(1.35),
      head_color=TEAL, body_fill=TEAL_LIGHT, size=22, body_align=PP_ALIGN.CENTER)
delta = rounded(s, Inches(4.4), Inches(3.65), Inches(4.5), Inches(0.6),
                UPM_BLUE, radius=0.5)
delta.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
rich(delta.text_frame.paragraphs[0], "Δ = −0.048  of fake performance",
     size=17, color=WHITE, base_bold=True, align=PP_ALIGN.CENTER)
bullets(s, [
    "That gap is the difference between a reproducible result and a broken one.",
    "A human reviewer might catch this. **At scale, nobody reviews every "
    "feature of every dataset.**",
    "The framework presented today detects this automatically — and "
    "quantifies the damage.",
], Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.3), size=16, gap=10)

# ======================================================================
# 5 — The gap
# ======================================================================
s = new_slide()
chrome(s, "The Gap in Existing Tools", kicker="Motivation & Objectives")
tools = [
    ("ydata-profiling", "Descriptive profiling, correlation reports",
     "0 / 4 leakage scenarios"),
    ("Great Expectations", "Rule-based validation suites",
     "1 / 4 (manual rule only)"),
    ("Deepchecks", "Train/test validation checks",
     "0 / 4 leakage scenarios"),
]
tx = Inches(0.7)
for name, desc, det in tools:
    card = rounded(s, tx, Inches(1.4), Inches(3.9), Inches(2.1), BLUE_PALE,
                   line=GRAY_LIGHT, line_w=1)
    box = textbox(s, tx + Inches(0.2), Inches(1.6), Inches(3.5), Inches(1.8))
    tf = box.text_frame
    rich(tf.paragraphs[0], name, size=16, color=UPM_BLUE, base_bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(5)
    rich(p, desc, size=12.5, color=DARK)
    p = tf.add_paragraph()
    p.space_before = Pt(9)
    rich(p, "Leakage detection:", size=11, color=GRAY_TXT)
    p = tf.add_paragraph()
    rich(p, "**" + det + "**", size=13, color=RED)
    tx += Inches(4.05)
bullets(s, [
    "Strong at profiling and rule-based validation — **but leakage detection is "
    "almost absent**, and none produces an interpretable readiness verdict or "
    "code-level fixes.",
], Inches(0.7), Inches(3.8), Inches(11.9), Inches(0.9), size=16, gap=6)
block(s, "Positioning",
      "The framework **complements** these tools — it fills the "
      "leakage-detection and actionability gap rather than replacing "
      "general-purpose profiling.",
      Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.4), size=16)

# ======================================================================
# 6 — Objectives & contributions
# ======================================================================
s = new_slide()
chrome(s, "Objectives and Contributions", kicker="Motivation & Objectives")
cw, ch = Inches(5.85), Inches(1.75)
numbered_card(s, 1, "A unified, automated framework",
              "20 deterministic checks across 5 dimensions (quality, leakage, "
              "features, sufficiency, drift) + impact analysis — via CLI, "
              "Python SDK, and web app.",
              Inches(0.7), Inches(1.3), cw, ch)
numbered_card(s, 2, "A unified Leakage Risk Score L(f)",
              "Correlation, mutual information, and single-feature performance "
              "inflation combined into one flaggable number per feature.",
              Inches(6.75), Inches(1.3), cw, ch)
numbered_card(s, 3, "LLM-based semantic leakage analysis",
              "Detects leakage from feature meaning (e.g. discharge_date) "
              "where statistics alone are blind.",
              Inches(0.7), Inches(3.2), cw, ch)
numbered_card(s, 4, "A quantitative benchmark",
              "Against three widely used tools on a 29-item checklist and four "
              "constructed leakage scenarios.",
              Inches(6.75), Inches(3.2), cw, ch)
kx = Inches(0.95)
for value, label in [("20", "automated checks,\n5 dimensions"),
                     ("12", "datasets\n(6 synthetic + 6 real)"),
                     ("4/4", "leakage scenarios\ndetected"),
                     ("325", "unit tests\npassing")]:
    kpi_card(s, value, label.replace("\n", " "), kx, Inches(5.3), w=Inches(2.7),
             h=Inches(1.35))
    kx += Inches(2.92)

# ======================================================================
# DIVIDER 02 — System & Methodology
# ======================================================================
divider(1, [
    "Pipeline architecture and the three user interfaces",
    "The 20-check catalog  ·  unified Leakage Risk Score with ablation",
    "LLM semantic analysis  ·  from 20 checks to one readiness score",
])

# ======================================================================
# 7 — Architecture
# ======================================================================
s = new_slide()
chrome(s, "System Architecture", kicker="System & Methodology")
BW, BH = Inches(2.15), Inches(0.62)


def cxc(shape):
    return shape.left + shape.width // 2


cfg = box_node(s, "config.yaml", Inches(4.0), Inches(1.25), BW, Inches(0.5),
               ORANGE_LIGHT, UPM_ORANGE)
dat = box_node(s, "Dataset CSV", Inches(7.2), Inches(1.25), BW, Inches(0.5),
               ORANGE_LIGHT, UPM_ORANGE)
load = box_node(s, "load_dataset()", Inches(5.6), Inches(2.05), BW, Inches(0.5),
                BLUE_LIGHT, UPM_BLUE)
row_y = Inches(2.95)
qual = box_node(s, "quality checks (6)", Inches(1.15), row_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
leak = box_node(s, "leakage checks (5)", Inches(4.15), row_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
feat = box_node(s, "feature analysis (3)", Inches(7.15), row_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
suff = box_node(s, "sufficiency checks (4)", Inches(10.15), row_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
row2_y = Inches(4.05)
drift = box_node(s, "drift detection (2)", Inches(2.65), row2_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
impact = box_node(s, "impact analysis", Inches(5.6), row2_y, BW, BH, BLUE_LIGHT, UPM_BLUE)
sem = box_node(s, "semantic leakage (optional)", Inches(8.55), row2_y, BW, BH,
               BLUE_LIGHT, UPM_BLUE)
rec = box_node(s, "recommendations + readiness score", Inches(4.1), Inches(5.2),
               Inches(2.5), BH, TEAL_LIGHT, TEAL)
rep = box_node(s, "HTML report + JSON export", Inches(7.3), Inches(5.2),
               Inches(2.5), BH, TEAL_LIGHT, TEAL)

arrow(s, cxc(cfg), cfg.top + cfg.height, load.left + Inches(0.4), load.top)
arrow(s, cxc(dat), dat.top + dat.height, load.left + load.width - Inches(0.4), load.top)
for tgt in (qual, leak, feat, suff):
    arrow(s, cxc(load), load.top + load.height, cxc(tgt), tgt.top)
arrow(s, cxc(qual), qual.top + qual.height, drift.left + Inches(0.3), drift.top)
arrow(s, cxc(leak), leak.top + leak.height, cxc(drift) + Inches(0.3), drift.top)
arrow(s, cxc(feat), feat.top + feat.height, cxc(impact) + Inches(0.3), impact.top)
arrow(s, cxc(suff), suff.top + suff.height, sem.left + sem.width - Inches(0.3), sem.top)
arrow(s, cxc(drift), drift.top + drift.height, rec.left + Inches(0.4), rec.top)
arrow(s, cxc(impact), impact.top + impact.height, cxc(rec) + Inches(0.2), rec.top)
arrow(s, cxc(sem), sem.top + sem.height, rec.left + rec.width - Inches(0.2), rec.top)
arrow(s, rec.left + rec.width, rec.top + rec.height // 2,
      rep.left, rep.top + rep.height // 2)
note(s, "One YAML config + one CSV in  →  20 checks in 5 dimensions  →  "
        "recommendations, readiness score, and reports out.", y=Inches(6.35))

# ======================================================================
# 8 — Interfaces
# ======================================================================
s = new_slide()
chrome(s, "Three Interfaces, One Pipeline", kicker="System & Methodology")
iface = [
    ("CLI", "main.py — batch runs, CI/CD friendly"),
    ("Python SDK", "DatasetChecker — notebooks and pipelines"),
    ("Streamlit web app", "seven-tab dashboard for non-programmers"),
]
iy = Inches(1.4)
for name, desc in iface:
    card = rounded(s, Inches(0.7), iy, Inches(5.4), Inches(1.15), BLUE_PALE,
                   line=GRAY_LIGHT, line_w=1)
    box = textbox(s, Inches(0.95), iy + Inches(0.13), Inches(4.9), Inches(0.95))
    tf = box.text_frame
    rich(tf.paragraphs[0], name, size=16, color=UPM_BLUE, base_bold=True)
    p = tf.add_paragraph()
    rich(p, desc, size=13, color=DARK)
    iy += Inches(1.33)
note(s, "Plus a @register_check **plugin API** for custom, domain-specific "
        "checks.", x=Inches(0.7), y=Inches(5.5), w=Inches(5.4), size=13)
code_box(s, [
    "from src.checker import DatasetChecker",
    "",
    "checker = DatasetChecker(",
    '    "configs/config.yaml")',
    "report = checker.run(",
    '    "data/titanic.csv",',
    '    target_col="survived")',
    "",
    "print(checker.score, checker.grade)",
    'checker.save_report("reports/")',
], Inches(6.6), Inches(1.4), Inches(6.1), Inches(3.9), size=14)

# ======================================================================
# 9 — Check catalog
# ======================================================================
s = new_slide()
chrome(s, "The 20-Check Catalog (5 Dimensions)", kicker="System & Methodology")
styled_table(s, [
    ["Dimension", "#", "Checks"],
    ["Quality", "6", "missing values, duplicates, outliers, class imbalance, "
     "constant features, low variance"],
    ["Leakage", "5", "target leakage, train/test overlap, temporal, ID columns, "
     "**unified risk score**"],
    ["Features", "3", "correlation structure, MI relevance, distribution shape"],
    ["Sufficiency", "4", "sample size, n/p ratio, class support, power"],
    ["Drift", "2", "KS test, PSI (covariate + label)"],
], Inches(0.9), Inches(1.3), Inches(11.5), row_h=0.55, size=14,
    col_widths=[2.0, 0.8, 8.7], align_center_cols=(1,))
bullets(s, [
    "Each check returns a structured CheckResult (pass/fail, severity, "
    "affected columns).",
    "Failed checks map to **20 recommendation handlers** producing "
    "code-level fixes.",
    "Optional 21st check: LLM semantic leakage (disabled by default).",
], Inches(0.9), Inches(4.9), Inches(11.5), Inches(2.0), size=16, gap=9)

# ======================================================================
# 10 — LRS
# ======================================================================
s = new_slide()
chrome(s, "Unified Leakage Risk Score  L(f)", kicker="System & Methodology")
eq = rounded(s, Inches(2.4), Inches(1.3), Inches(8.5), Inches(0.85), BLUE_LIGHT,
             line=UPM_BLUE, line_w=1, radius=0.3)
eq.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
rich(eq.text_frame.paragraphs[0],
     "L(f)  =  0.35 · ρ(f)  +  0.35 · Ĩ(f; y)  +  0.30 · π(f)",
     size=23, base_bold=True, align=PP_ALIGN.CENTER, color=UPM_BLUE)
signals = [
    ("ρ(f) — correlation", "Pearson |r| (numeric) or Cramér's V (categorical), "
     "in [0, 1]"),
    ("Ĩ(f; y) — mutual information", "normalized by the maximum MI across "
     "features, in [0, 1]"),
    ("π(f) — performance inflation", "(A_f − A_base) / (1 − A_base): how close "
     "does this single feature get a model to perfect accuracy?"),
]
sx = Inches(0.7)
for name, desc in signals:
    card = rounded(s, sx, Inches(2.55), Inches(3.9), Inches(1.9), BLUE_PALE,
                   line=GRAY_LIGHT, line_w=1)
    box = textbox(s, sx + Inches(0.2), Inches(2.75), Inches(3.5), Inches(1.6))
    tf = box.text_frame
    rich(tf.paragraphs[0], name, size=14.5, color=UPM_BLUE, base_bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(5)
    rich(p, desc, size=12.5, color=DARK)
    sx += Inches(4.05)
block(s, "Flagging thresholds",
      "L(f) ≥ 0.7  ⇒  **warning**             L(f) ≥ 0.9  ⇒  **error**",
      Inches(2.4), Inches(4.95), Inches(8.5), Inches(1.15), size=18,
      body_align=PP_ALIGN.CENTER)
note(s, "Single statistics miss real leaks; three complementary signals are "
        "far more robust.", y=Inches(6.4))

# ======================================================================
# 11 — Ablation
# ======================================================================
s = new_slide()
chrome(s, "Are the Weights Arbitrary? — Ablation", kicker="System & Methodology")
intro = textbox(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.45))
rich(intro.text_frame.paragraphs[0],
     "Weight-sensitivity analysis on known leaky features (Titanic):", size=17)
styled_table(s, [
    ["Weight scheme  (ρ / Ĩ / π)", "L(name)", "L(boat)"],
    ["**Default (0.35 / 0.35 / 0.30)**", "**1.000**", "**0.74**  ✓ flagged"],
    ["Correlation-heavy (0.50 / 0.25 / 0.25)", "1.000",
     "0.662  ✗ below threshold!"],
    ["MI- or π-weighted schemes", "1.000", "flagged  ✓"],
], Inches(1.6), Inches(1.8), Inches(10.1), row_h=0.52, size=15,
    col_widths=[4.6, 2.2, 3.3], align_center_cols=(1, 2))
bullets(s, [
    "boat (lifeboat number) is a true leak, but its raw correlation is diluted "
    "by missing values — an over-weighted correlation term produces a "
    "**false negative**.",
    "Every scheme that gives reasonable weight to MI or performance inflation "
    "catches it.",
    "Empirical evidence that the default weights are **defensible, not "
    "decorative**.",
], Inches(0.7), Inches(4.35), Inches(11.9), Inches(2.5), size=16, gap=11)

# ======================================================================
# 12 — Semantic LLM
# ======================================================================
s = new_slide()
chrome(s, "Semantic Leakage Analysis with an LLM", kicker="System & Methodology")
head = textbox(s, Inches(0.7), Inches(1.25), Inches(6.4), Inches(0.4))
rich(head.text_frame.paragraphs[0], "Why statistics are not enough", size=17,
     color=UPM_BLUE, base_bold=True)
bullets(s, [
    "A feature named discharge_date in an ICU mortality dataset is leaky **by "
    "meaning**: it may show weak correlation, yet it cannot exist at "
    "prediction time.",
    "No statistical test reads feature **names**.",
], Inches(0.7), Inches(1.7), Inches(6.4), Inches(2.0), size=15, gap=8)
head = textbox(s, Inches(0.7), Inches(3.85), Inches(6.4), Inches(0.4))
rich(head.text_frame.paragraphs[0], "Approach", size=17, color=UPM_BLUE,
     base_bold=True)
bullets(s, [
    "Feature names + sample values + dataset description → **GPT-4o-mini** "
    "(Azure OpenAI).",
    "Degrades gracefully without credentials; disabled by default.",
], Inches(0.7), Inches(4.3), Inches(6.4), Inches(2.0), size=15, gap=8)
block(s, "Per-feature output", [
    "risk_level:   none / low / medium / high",
    "leakage_type:   temporal / proxy / post-hoc / indirect",
    "+ natural-language rationale",
], Inches(7.5), Inches(1.4), Inches(5.2), Inches(2.15), size=14)
kpi_card(s, "0% → F1 0.963", "naming-based leakage detection, 30-feature "
         "labelled benchmark", Inches(7.5), Inches(3.9), w=Inches(5.2),
         h=Inches(1.5), accent=TEAL)

# ======================================================================
# 13 — Semantic caveat
# ======================================================================
s = new_slide()
chrome(s, "Semantic Module: Evaluation and an Honest Caveat",
       kicker="System & Methodology")
bullets(s, [
    "Manually labelled **30-feature benchmark** across temporal, proxy, "
    "post-hoc, and indirect leakage types (data/semantic_benchmark.json).",
], Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.9), size=17, gap=8)
kx = Inches(2.4)
for value, label in [("1.00", "Precision"), ("0.93", "Recall"), ("0.963", "F1")]:
    kpi_card(s, value, label, kx, Inches(2.15), w=Inches(2.6), h=Inches(1.3))
    kx += Inches(2.95)
block(s, "Transparency",
      "These figures validate the **evaluation harness** using a deterministic "
      "mock analyzer — not the live GPT-4o-mini model (Azure credentials were "
      "unavailable at evaluation time). A live-model evaluation is the "
      "**top-priority item of future work**, and this is stated explicitly in "
      "the thesis.",
      Inches(0.7), Inches(3.85), Inches(11.9), Inches(1.85),
      head_color=RED, body_fill=RED_LIGHT, size=15)
note(s, "The architecture, prompt design, structured output parsing, and "
        "graceful degradation are fully implemented and tested "
        "(33 unit tests for Phase 16 modules).", y=Inches(6.0))

# ======================================================================
# 14 — Readiness score
# ======================================================================
s = new_slide()
chrome(s, "From 20 Checks to One Number: the Readiness Score",
       kicker="System & Methodology")
eq = rounded(s, Inches(1.5), Inches(1.35), Inches(10.3), Inches(0.85),
             BLUE_LIGHT, line=UPM_BLUE, line_w=1, radius=0.3)
eq.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
rich(eq.text_frame.paragraphs[0],
     "overall  =  0.25·S_quality  +  0.35·S_leakage  +  0.25·S_features  +  "
     "0.15·S_sufficiency",
     size=19, base_bold=True, align=PP_ALIGN.CENTER, color=UPM_BLUE)
bullets(s, [
    "Per dimension: start at 100, subtract **15 per error**, **5 per warning**.",
    "**Leakage carries the largest weight (0.35)**: it is the only failure mode "
    "that silently inflates results instead of visibly degrading them.",
    "All weights and penalties are **configurable** in config.yaml.",
], Inches(0.9), Inches(2.6), Inches(11.6), Inches(2.3), size=17, gap=12)
grades = [("A ≥ 85", GREEN_LIGHT), ("B ≥ 70", TEAL_LIGHT),
          ("C ≥ 55", RGBColor(0xFF, 0xE9, 0xCC)),
          ("D ≥ 40", RGBColor(0xFF, 0xD9, 0xA8)),
          ("F < 40", RGBColor(0xF6, 0xD0, 0xD0))]
gx = Inches(2.05)
for label, color in grades:
    chip = rounded(s, gx, Inches(5.35), Inches(1.65), Inches(0.6), color,
                   line=GRAY, line_w=0.5, radius=0.35)
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(chip.text_frame.paragraphs[0], label, size=16, base_bold=True,
         align=PP_ALIGN.CENTER)
    gx += Inches(1.9)

# ======================================================================
# DIVIDER 03 — Experimental Results
# ======================================================================
divider(2, [
    "Twelve datasets: six synthetic with planted defects, six real-world",
    "Readiness scores  ·  the Titanic case study  ·  quantifying leakage damage",
])

# ======================================================================
# 15 — Datasets
# ======================================================================
s = new_slide()
chrome(s, "Experimental Setup: 12 Datasets", kicker="Experimental Results")
head = textbox(s, Inches(0.7), Inches(1.25), Inches(5.9), Inches(0.4))
rich(head.text_frame.paragraphs[0], "6 synthetic — controlled ground truth",
     size=16, color=UPM_BLUE, base_bold=True)
card = rounded(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(3.35), BLUE_PALE,
               line=GRAY_LIGHT, line_w=1)
bullets(s, [
    "clean_dataset — control",
    "dirty_dataset — quality issues",
    "leaky_dataset — leakage issues",
    "proxy_leakage — graded noisy proxies",
    "temporal_leakage_ext — future feature",
    "multitype_leakage — ICU: proxy + temporal + ID",
], Inches(0.95), Inches(1.9), Inches(5.5), Inches(3.0), size=14, gap=7)
head = textbox(s, Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.4))
rich(head.text_frame.paragraphs[0], "6 real-world — OpenML / UCI", size=16,
     color=UPM_BLUE, base_bold=True)
card = rounded(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(3.35), BLUE_PALE,
               line=GRAY_LIGHT, line_w=1)
bullets(s, [
    "Titanic (1,309 rows)",
    "Pima Diabetes (768)",
    "Adult Census (48,842)",
    "German Credit (1,000)",
    "Heart Disease (303)",
    "Wine Quality (1,599)",
], Inches(7.15), Inches(1.9), Inches(5.5), Inches(3.0), size=14, gap=7)
bullets(s, [
    "Synthetic sets verify the framework **detects what we planted**; real "
    "sets verify it **finds surprises in the wild**.",
    "Full pipeline per dataset: 20 checks + impact analysis (23 results).",
], Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.5), size=16, gap=9)

# ======================================================================
# 16 — Readiness chart
# ======================================================================
s = new_slide()
chrome(s, "Readiness Scores Discriminate as Designed",
       kicker="Experimental Results")
chart_data = CategoryChartData()
chart_data.categories = ["clean_dataset", "dirty_dataset", "leaky_dataset",
                         "Titanic", "Pima Diabetes"]
chart_data.add_series("Readiness Score", [98.8, 91.2, 71.6, 80.6, 96.2])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.1),
                        Inches(1.25), Inches(11.1), Inches(4.4), chart_data)
chart = gf.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 80
plot.has_data_labels = True
plot.data_labels.font.size = Pt(14)
plot.data_labels.font.bold = True
plot.data_labels.number_format = "0.0"
plot.data_labels.number_format_is_linked = False
series = plot.series[0]
for i, color in enumerate([TEAL, TEAL, UPM_ORANGE, UPM_ORANGE, TEAL]):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color
va = chart.value_axis
va.maximum_scale = 110
va.minimum_scale = 0
va.major_unit = 20
va.tick_labels.font.size = Pt(12)
chart.category_axis.tick_labels.font.size = Pt(13)
legend = textbox(s, Inches(1.1), Inches(5.7), Inches(11.1), Inches(0.4))
par = legend.text_frame.paragraphs[0]
par.alignment = PP_ALIGN.CENTER
for txt, color in [("■ Grade A (≥ 85)        ", TEAL),
                   ("■ Grade B (≥ 70)", UPM_ORANGE)]:
    run = par.add_run()
    run.text = txt
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color
note(s, "Only the datasets with **confirmed target leakage** drop to grade B — "
        "quality noise alone (dirty) stays A.  UCI suite: Adult 92.4, German "
        "Credit 97.5, Heart Disease 96.8, Wine Quality 88.6 (all A).",
     y=Inches(6.25))

# ======================================================================
# 17 — Titanic case study
# ======================================================================
s = new_slide()
chrome(s, "Case Study: Titanic — Real Leaks in a Classic Dataset",
       kicker="Experimental Results")
feats = [
    ("name", "Cramér's V = 0.999 with survived", "target leakage — ERROR",
     "Unique names act as row identifiers.", RED),
    ("boat", "Unified risk score L = 0.74", "post-hoc leakage — WARNING",
     "Lifeboat number is assigned after the outcome: knowing it literally "
     "encodes survival.", UPM_ORANGE),
]
fy = Inches(1.35)
for fname, stat, verdict, expl, color in feats:
    card = rounded(s, Inches(0.7), fy, Inches(7.0), Inches(1.85), BLUE_PALE,
                   line=GRAY_LIGHT, line_w=1)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), fy,
                               Inches(0.07), Inches(1.85))
    set_fill(strip, color)
    box = textbox(s, Inches(0.95), fy + Inches(0.12), Inches(6.6), Inches(1.65))
    tf = box.text_frame
    par = tf.paragraphs[0]
    run = par.add_run()
    run.text = fname
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = "Consolas"
    run.font.color.rgb = UPM_BLUE
    run = par.add_run()
    run.text = "    " + stat
    run.font.size = Pt(13)
    run.font.color.rgb = DARK
    p = tf.add_paragraph()
    rich(p, "**" + verdict + "**", size=13, color=color)
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    rich(p, expl, size=12.5, color=DARK)
    fy += Inches(2.05)
bullets(s, [
    "Neither is caught by any of the three baseline tools.",
], Inches(0.7), Inches(5.5), Inches(7.0), Inches(0.6), size=15, gap=6)
block(s, "Why this matters",
      "Titanic is used in thousands of tutorials and courses. Models "
      "“achieving” high accuracy with boat included are **textbook post-hoc "
      "leakage** — and it takes an automated tool 30 seconds to prove it.",
      Inches(8.0), Inches(1.35), Inches(4.7), Inches(2.6), size=13.5)
kpi_card(s, "80.6 / B", "Titanic readiness (14/22 checks passed)",
         Inches(8.0), Inches(4.3), w=Inches(4.7), h=Inches(1.35))
note(s, "Additional case studies (thesis §5.5): Adult Census (missing values, "
        "skew), German Credit (near-zero-MI features), Wine Quality "
        "(15% duplicates → error).", y=Inches(6.45))

# ======================================================================
# 18 — Impact analysis
# ======================================================================
s = new_slide()
chrome(s, "Quantifying the Damage: Impact Analysis",
       kicker="Experimental Results")
intro = textbox(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.45))
rich(intro.text_frame.paragraphs[0],
     "The pipeline retrains the same model with and without the flagged "
     "features (cross-validated):", size=17)
styled_table(s, [
    ["Dataset (LR)", "Baseline acc.", "Cleaned acc.", "Δ"],
    ["leaky_dataset", "1.000", "0.952", "**−0.048**"],
], Inches(2.6), Inches(1.85), Inches(8.1), row_h=0.55, size=16,
    col_widths=[2.7, 2.0, 2.0, 1.4], align_center_cols=(1, 2, 3))
bullets(s, [
    "A **positive** baseline−cleaned gap is direct, quantified evidence of "
    "leakage — not just a statistical suspicion.",
    "This closes the loop on the motivating example: the framework **finds** "
    "the leak, **explains** it, and **measures** how much of the reported "
    "performance was fake.",
    "Recommendations then emit the exact pandas code to drop or re-derive "
    "the offending columns.",
], Inches(0.7), Inches(3.45), Inches(11.9), Inches(3.0), size=17, gap=12)

# ======================================================================
# DIVIDER 04 — Benchmark
# ======================================================================
divider(3, [
    "Feature coverage on a 29-item cross-tool checklist",
    "Leakage detection on four constructed scenarios — the decisive result",
])

# ======================================================================
# 19 — Coverage chart
# ======================================================================
s = new_slide()
chrome(s, "Feature Coverage vs. Three Established Tools", kicker="Benchmark")
chart_data = CategoryChartData()
chart_data.categories = ["Great Expectations", "Deepchecks", "ydata-profiling",
                         "ml-framework"]
chart_data.add_series("Items satisfied (of 29)", [9, 11, 8, 29])
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.4), Inches(1.25),
                        Inches(10.5), Inches(3.6), chart_data)
chart = gf.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 70
plot.has_data_labels = True
plot.data_labels.font.size = Pt(14)
plot.data_labels.font.bold = True
series = plot.series[0]
for i, color in enumerate([RGBColor(0xC8, 0xDC, 0xF0), RGBColor(0xB4, 0xD2, 0xE6),
                           RGBColor(0xA0, 0xC3, 0xE1), TEAL]):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color
va = chart.value_axis
va.maximum_scale = 30
va.minimum_scale = 0
va.major_unit = 5
va.tick_labels.font.size = Pt(12)
chart.category_axis.tick_labels.font.size = Pt(13)
bullets(s, [
    "The 29-item checklist is a **cross-tool comparison instrument** (more "
    "granular than the 20-check pipeline; the two numbers measure different "
    "things).",
    "**Fairness note**: the baselines target profiling/validation, not leakage "
    "— the comparison measures scope, and the checklist was published with "
    "the benchmark code for scrutiny.",
], Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.9), size=14.5, gap=8)

# ======================================================================
# 20 — Detection table
# ======================================================================
s = new_slide()
chrome(s, "Leakage Detection — the Decisive Result", kicker="Benchmark")
intro = textbox(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.45))
rich(intro.text_frame.paragraphs[0],
     "Four constructed leakage scenarios, each tool run with default "
     "configuration:", size=17)
styled_table(s, [
    ["Scenario", "ml-framework", "ydata", "Deepchecks", "Great Exp."],
    ["Perfect proxy (r = 1.0)", "✓", "✗", "✗", "✗"],
    ["Noisy proxy (r ≈ 0.97)", "✓", "✗", "✗", "✗"],
    ["ID column (100% cardinality)", "✓", "✗", "✗", "✓"],
    ["Indirect computed feature", "✓", "✗", "✗", "✗"],
    ["**Detection rate**", "**4/4**", "0/4", "0/4", "1/4"],
], Inches(1.3), Inches(1.8), Inches(10.7), row_h=0.5, size=15,
    col_widths=[3.9, 2.0, 1.5, 1.8, 1.5], align_center_cols=(1, 2, 3, 4))
bullets(s, [
    "Great Expectations' single detection required a **manually authored "
    "rule**; ml-framework's four are fully automatic.",
    "This is precisely the gap the thesis set out to fill.",
], Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.7), size=17, gap=10)

# ======================================================================
# DIVIDER 05 — Demo & Conclusions
# ======================================================================
divider(4, [
    "Live Streamlit demo on the Titanic dataset",
    "Conclusions  ·  future research directions",
])

# ======================================================================
# 21 — Demo
# ======================================================================
s = new_slide()
chrome(s, "Live Demo: Streamlit Web Application", kicker="Demo & Conclusions")
head = textbox(s, Inches(0.7), Inches(1.25), Inches(6.6), Inches(0.4))
rich(head.text_frame.paragraphs[0], "What you will see (titanic.csv):",
     size=17, color=UPM_BLUE, base_bold=True)
steps = [
    "Upload the CSV, select survived as target.",
    "Pipeline runs live: 20 checks + impact analysis.",
    "Readiness score **80.6 / B** with per-dimension breakdown.",
    "name and boat flagged with explanations and suggested fixes.",
    "Downloadable HTML report.",
]
sy = Inches(1.85)
for i, step in enumerate(steps):
    num_circle(s, i + 1, Inches(0.75), sy, d=Inches(0.42), size=14)
    box = textbox(s, Inches(1.35), sy + Inches(0.015), Inches(5.9), Inches(0.75))
    rich(box.text_frame.paragraphs[0], step, size=14.5)
    sy += Inches(0.88)
block(s, "Run it yourself", ["streamlit run app/app.py"],
      Inches(7.7), Inches(1.4), Inches(5.0), Inches(1.1), size=16,
      body_align=PP_ALIGN.CENTER)
box = textbox(s, Inches(7.7), Inches(2.85), Inches(5.0), Inches(1.8))
tf = box.text_frame
rich(tf.paragraphs[0],
     "Seven tabs: overview, quality, leakage, features, sufficiency, impact, "
     "readiness summary.", size=13, color=GRAY_TXT)
p = tf.add_paragraph()
p.space_before = Pt(10)
rich(p, "(Backup: pre-recorded capture and generated HTML report, in case "
        "of demo issues.)", size=13, color=GRAY_TXT)

# ======================================================================
# 22 — Conclusions
# ======================================================================
s = new_slide()
chrome(s, "Conclusions", kicker="Demo & Conclusions")
cw, ch = Inches(5.85), Inches(1.8)
numbered_card(s, 1, "Automation is feasible end-to-end",
              "One config + one CSV → scored, explained, actionable verdict on "
              "dataset quality and leakage.",
              Inches(0.7), Inches(1.3), cw, ch, accent=TEAL)
numbered_card(s, 2, "The unified L(f) score works",
              "Detects leaks no single statistic catches, with ablation-backed "
              "weights — 4/4 benchmark scenarios vs. 0–1/4 for existing tools.",
              Inches(6.75), Inches(1.3), cw, ch, accent=TEAL)
numbered_card(s, 3, "LLMs extend detection to semantics",
              "Leakage invisible to statistics becomes detectable; evaluation "
              "harness built and validated, live-model run queued.",
              Inches(0.7), Inches(3.25), cw, ch, accent=TEAL)
numbered_card(s, 4, "Real issues found in classic datasets",
              "Titanic boat/name and Wine Quality duplicates — previously "
              "unflagged by standard tooling.",
              Inches(6.75), Inches(3.25), cw, ch, accent=TEAL)
kx = Inches(1.65)
for value, label in [("17", "development phases"), ("14", "source modules"),
                     ("325/325", "tests passing"), ("3", "user interfaces")]:
    kpi_card(s, value, label, kx, Inches(5.35), w=Inches(2.4), h=Inches(1.25))
    kx += Inches(2.6)

# ======================================================================
# 23 — Future work
# ======================================================================
s = new_slide()
chrome(s, "Future Work", kicker="Demo & Conclusions")
fw = [
    ("Live LLM evaluation", "of the semantic module against the 30-feature "
     "benchmark — top priority; the harness is ready, only credentials are "
     "missing.", UPM_ORANGE),
    ("Broader leakage taxonomy", "group leakage, preprocessing leakage (scaler "
     "fit on full data), cross-validation contamination.", UPM_BLUE),
    ("CI/CD integration", "readiness score as a merge gate for data pipelines "
     "— the SDK already makes this a few lines of code.", UPM_BLUE),
    ("Learning thresholds from data", "replacing fixed flagging thresholds "
     "with calibrated, dataset-size-aware values.", UPM_BLUE),
]
fy = Inches(1.4)
for i, (title, desc, accent) in enumerate(fw):
    numbered_card(s, i + 1, title, desc, Inches(1.3), fy, Inches(10.7),
                  Inches(1.25), accent=accent)
    fy += Inches(1.42)

# ======================================================================
# 24 — Thank you (full-bleed)
# ======================================================================
_slide_no += 1
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
set_fill(bg, UPM_BLUE)
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
set_fill(side, UPM_ORANGE)
t = textbox(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.2))
rich(t.text_frame.paragraphs[0], "Thank you", size=54, color=WHITE,
     base_bold=True, align=PP_ALIGN.CENTER)
q = textbox(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(0.6))
rich(q.text_frame.paragraphs[0], "Questions?", size=22, color=BLUE_SOFT,
     align=PP_ALIGN.CENTER)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.07), Inches(4.25),
                          Inches(3.2), Inches(0.035))
set_fill(rule, UPM_ORANGE)
info = textbox(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(1.6))
tf = info.text_frame
rich(tf.paragraphs[0], "github.com/jaimecano12/ml-framework", size=18,
     color=WHITE, base_bold=True, align=PP_ALIGN.CENTER)
p = tf.add_paragraph()
p.space_before = Pt(10)
rich(p, "Jaime Cano Moraño", size=16, color=WHITE, align=PP_ALIGN.CENTER)
p = tf.add_paragraph()
rich(p, "ETSIT-UPM  ·  Illinois Institute of Technology", size=13,
     color=BLUE_SOFT, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------------
prs.save("presentation.pptx")
print(f"Saved presentation.pptx with {len(prs.slides._sldIdLst)} slides")
